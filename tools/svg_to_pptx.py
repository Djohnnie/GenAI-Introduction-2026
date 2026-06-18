"""
svg_to_pptx.py — Convert a folder of SVG slide files into a PowerPoint file.

Usage:
    python svg_to_pptx.py <svg_dir> [output.pptx]

Arguments:
    svg_dir      Directory containing slide-*.svg files (e.g. content/2-day/chapter-01)
    output.pptx  Optional output path. Defaults to <svg_dir>/<dirname>.pptx

Requirements:
    pip install python-pptx
    Microsoft Edge installed (used for headless SVG-to-PNG rendering)
"""

import argparse
import glob
import os
import subprocess
import sys
import tempfile

from pptx import Presentation
from pptx.util import Emu

EDGE_PATHS = [
    r"C:\Program Files (x86)\Microsoft\Edge\Application\msedge.exe",
    r"C:\Program Files\Microsoft\Edge\Application\msedge.exe",
]

# Standard 16:9 widescreen slide dimensions in EMU (English Metric Units)
SLIDE_WIDTH  = Emu(12192000)  # 13.333 inches
SLIDE_HEIGHT = Emu(6858000)   # 7.5 inches

# Edge renders at this resolution to match the SVG viewBox (1920×1080)
RENDER_WIDTH  = 1920
RENDER_HEIGHT = 1080


def find_edge() -> str:
    for path in EDGE_PATHS:
        if os.path.isfile(path):
            return path
    raise FileNotFoundError(
        "Microsoft Edge not found. Install Edge or adjust EDGE_PATHS in this script."
    )


def svg_to_png(edge: str, svg_path: str, png_path: str, timeout: int = 30) -> bool:
    url = "file:///" + svg_path.replace("\\", "/")
    result = subprocess.run(
        [edge, "--headless=new", "--disable-gpu",
         f"--screenshot={png_path}",
         f"--window-size={RENDER_WIDTH},{RENDER_HEIGHT}",
         url],
        capture_output=True,
        timeout=timeout,
    )
    return os.path.isfile(png_path)


def build_pptx(svg_dir: str, output_path: str) -> None:
    svg_dir = os.path.abspath(svg_dir)
    svgs = sorted(glob.glob(os.path.join(svg_dir, "slide-*.svg")))

    if not svgs:
        print(f"No slide-*.svg files found in: {svg_dir}", file=sys.stderr)
        sys.exit(1)

    print(f"Found {len(svgs)} SVG(s) in {svg_dir}")

    edge = find_edge()

    prs = Presentation()
    prs.slide_width  = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT
    blank_layout = prs.slide_layouts[6]  # completely blank

    with tempfile.TemporaryDirectory() as tmp:
        for i, svg in enumerate(svgs, 1):
            name = os.path.splitext(os.path.basename(svg))[0]
            png  = os.path.join(tmp, f"{name}.png")

            ok = svg_to_png(edge, svg, png)
            if not ok:
                print(f"  [{i:02d}/{len(svgs)}] WARNING: could not render {name}, skipping")
                continue

            slide = prs.slides.add_slide(blank_layout)
            slide.shapes.add_picture(png, 0, 0, SLIDE_WIDTH, SLIDE_HEIGHT)
            print(f"  [{i:02d}/{len(svgs)}] {name}")

    prs.save(output_path)
    print(f"\nSaved: {output_path}")


def main() -> None:
    parser = argparse.ArgumentParser(
        description="Convert a folder of slide-*.svg files into a PowerPoint file."
    )
    parser.add_argument("svg_dir", help="Directory containing slide-*.svg files")
    parser.add_argument(
        "output",
        nargs="?",
        help="Output .pptx path (default: <svg_dir>/<dirname>.pptx)",
    )
    args = parser.parse_args()

    svg_dir = os.path.abspath(args.svg_dir)
    if not os.path.isdir(svg_dir):
        print(f"Directory not found: {svg_dir}", file=sys.stderr)
        sys.exit(1)

    output = args.output or os.path.join(svg_dir, os.path.basename(svg_dir) + ".pptx")
    output = os.path.abspath(output)

    build_pptx(svg_dir, output)


if __name__ == "__main__":
    main()
